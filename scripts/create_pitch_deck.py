"""虾塘大亨 · OPC极限挑战赛 路演PPT生成脚本。

基于观猹模板，生成 8 页路演 PPT。
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = "/Users/a647/.openclaw/media/inbound/PPT模版丨观猹丨OPC极限挑战赛---43a95559-b4df-486d-86c0-38ef8ad37e3b.pptx"
OUTPUT = "/Users/a647/projects/shrimp-tycoon/虾塘大亨_OPC路演.pptx"

# 配色
TEAL = RGBColor(0x00, 0xC8, 0xB4)       # 主色 #00c8b4
DARK_BG = RGBColor(0x0A, 0x0E, 0x1A)    # 深色背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC4)
ACCENT_ORANGE = RGBColor(0xF5, 0xA6, 0x23)
ALERT_RED = RGBColor(0xFF, 0x4D, 0x6D)
SEMI_WHITE = RGBColor(0xE0, 0xE4, 0xEB)

# 尺寸（基于模板 28" x 10.5"）
SLIDE_W = Inches(28)
SLIDE_H = Inches(10.5)


def set_slide_bg(slide, color):
    """设置幻灯片背景色。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, 
                 font_size=24, font_color=WHITE, bold=False, 
                 alignment=PP_ALIGN.LEFT, font_name="思源黑体 CN Bold"):
    """添加文本框。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=20, font_color=WHITE, 
                  bold=False, font_name="思源黑体 CN Bold",
                  space_before=Pt(6), alignment=PP_ALIGN.LEFT):
    """在已有文本框中追加段落。"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_rect(slide, left, top, width, height, fill_color, 
             border_color=None, corner_radius=None):
    """添加矩形。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def create_slide1_cover(prs):
    """Slide 1: 封面。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    
    # 项目名称
    add_text_box(slide, Inches(2), Inches(2.5), Inches(24), Inches(2),
                 "🦞 虾塘大亨", font_size=72, font_color=TEAL, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    # 副标题
    add_text_box(slide, Inches(2), Inches(4.5), Inches(24), Inches(1.5),
                 "AI 智慧水产决策系统", font_size=40, font_color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    # Slogan
    add_text_box(slide, Inches(2), Inches(6), Inches(24), Inches(1),
                 "不是让农民学 AI，是让 AI 学会养虾。", font_size=28, 
                 font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    # 赛道信息
    add_text_box(slide, Inches(2), Inches(7.5), Inches(24), Inches(1),
                 "OPC极限挑战赛 · 赛道二「AI合伙人」", font_size=22,
                 font_color=SEMI_WHITE, alignment=PP_ALIGN.CENTER)
    
    # 参赛者
    add_text_box(slide, Inches(2), Inches(8.5), Inches(24), Inches(0.8),
                 "参赛者：647", font_size=20, font_color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)


def create_slide2_problem(prs):
    """Slide 2: 行业痛点。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    
    # 标题
    add_text_box(slide, Inches(1.5), Inches(0.8), Inches(25), Inches(1.2),
                 "01  行业痛点：3,050 万亩虾塘，95% 在「盲养」", 
                 font_size=40, font_color=TEAL, bold=True)
    
    # 三个痛点卡片
    cards = [
        ("💀 死亡率高", "养殖全程无专业指导\n病害爆发后才发现\n年均损失率 30-50%",
         "传统：凭经验，靠天吃饭"),
        ("⏰ 人力依赖重", "每天 3-5 次人工巡塘\n夜间缺氧无人知\n单塘需 1-2 名技术员",
         "痛点：招不到、留不住、教不会"),
        ("📉 利润被挤压", "信息不对称卖不上价\n错过最佳捕捞窗口\n饲料投入产出比模糊",
         "现状：增产不增收"),
    ]
    
    for i, (title, desc, note) in enumerate(cards):
        x = Inches(1.5 + i * 8.5)
        # 卡片背景
        add_rect(slide, x, Inches(2.8), Inches(7.5), Inches(6.5), 
                 RGBColor(0x14, 0x1A, 0x2E), border_color=TEAL)
        
        # 卡片标题
        add_text_box(slide, x + Inches(0.5), Inches(3.2), Inches(6.5), Inches(1),
                     title, font_size=32, font_color=WHITE, bold=True)
        
        # 卡片内容
        add_text_box(slide, x + Inches(0.5), Inches(4.5), Inches(6.5), Inches(3),
                     desc, font_size=22, font_color=SEMI_WHITE)
        
        # 底部注释
        add_text_box(slide, x + Inches(0.5), Inches(7.8), Inches(6.5), Inches(0.8),
                     note, font_size=18, font_color=ACCENT_ORANGE, bold=True)


def create_slide3_solution(prs):
    """Slide 3: 解决方案。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    
    add_text_box(slide, Inches(1.5), Inches(0.8), Inches(25), Inches(1.2),
                 "02  解决方案：AI 合伙人，24/7 守护每一口塘", 
                 font_size=40, font_color=TEAL, bold=True)
    
    # 三个 Agent 卡片
    agents = [
        ("🛡️ 哨兵 Agent", "实时监控", 
         "每 5 分钟自动分析水质\n三层智能路由：规则→Haiku→Opus\n危险信号即时推送飞书\n10 秒内完成决策",
         "替代：24h值班技术员"),
        ("📊 策略 Agent", "每日决策",
         "自动生成日报（投喂/病害/捕捞）\n7 天趋势分析\n70 条知识库规则辅助推理\n精准投喂方案（减少30%饲料浪费）",
         "替代：养殖顾问"),
        ("📈 增长 Agent", "商业增长",
         "多塘 ROI 对比分析\n实时市场价格匹配\n最佳捕捞窗口预测\nICP 评分获客（客户画像）",
         "替代：销售经理+市场分析师"),
    ]
    
    for i, (title, subtitle, desc, replaces) in enumerate(agents):
        x = Inches(1.5 + i * 8.5)
        add_rect(slide, x, Inches(2.5), Inches(7.5), Inches(7), 
                 RGBColor(0x14, 0x1A, 0x2E), border_color=TEAL)
        
        add_text_box(slide, x + Inches(0.5), Inches(2.8), Inches(6.5), Inches(0.8),
                     title, font_size=32, font_color=WHITE, bold=True)
        
        add_text_box(slide, x + Inches(0.5), Inches(3.6), Inches(6.5), Inches(0.5),
                     subtitle, font_size=22, font_color=TEAL)
        
        add_text_box(slide, x + Inches(0.5), Inches(4.3), Inches(6.5), Inches(3.2),
                     desc, font_size=20, font_color=SEMI_WHITE)
        
        add_text_box(slide, x + Inches(0.5), Inches(7.8), Inches(6.5), Inches(0.8),
                     replaces, font_size=18, font_color=ACCENT_ORANGE, bold=True)


def create_slide4_architecture(prs):
    """Slide 4: 技术架构。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    
    add_text_box(slide, Inches(1.5), Inches(0.8), Inches(25), Inches(1.2),
                 "03  三层架构：决策层 → 工具层 → 数据层", 
                 font_size=40, font_color=TEAL, bold=True)
    
    # Layer A: Agent 决策层
    add_rect(slide, Inches(1.5), Inches(2.5), Inches(25), Inches(2.2), 
             RGBColor(0x14, 0x1A, 0x2E), border_color=TEAL)
    add_text_box(slide, Inches(2), Inches(2.7), Inches(4), Inches(0.6),
                 "Layer A · 决策层", font_size=24, font_color=TEAL, bold=True)
    add_text_box(slide, Inches(2), Inches(3.4), Inches(23), Inches(1),
                 "哨兵 Agent（实时）  →  策略 Agent（每日 cron）  →  增长 Agent（每周 cron）\n"
                 "三层路由：CSI≤20→规则引擎 | CSI 21-60→Haiku | CSI>60→Opus",
                 font_size=20, font_color=SEMI_WHITE)
    
    # Arrow
    add_text_box(slide, Inches(13), Inches(4.7), Inches(2), Inches(0.6),
                 "▼  MCP 协议", font_size=18, font_color=LIGHT_GRAY, 
                 alignment=PP_ALIGN.CENTER)
    
    # Layer B: MCP 工具层
    add_rect(slide, Inches(1.5), Inches(5.3), Inches(25), Inches(2.2),
             RGBColor(0x14, 0x1A, 0x2E), border_color=ACCENT_ORANGE)
    add_text_box(slide, Inches(2), Inches(5.5), Inches(4), Inches(0.6),
                 "Layer B · 工具层（12 个 MCP 工具）", font_size=24, 
                 font_color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(2), Inches(6.2), Inches(23), Inches(1),
                 "sensor_read · water_quality_score · feeding_recommend · disease_assess · harvest_advise\n"
                 "market_match · price_trend · kb_query · feishu_push · lead_score · crm_write · audit_log\n"
                 "特性：无状态 | 幂等 | ≤5s 超时 | 标准 FastMCP 协议",
                 font_size=18, font_color=SEMI_WHITE)
    
    # Arrow
    add_text_box(slide, Inches(13), Inches(7.5), Inches(2), Inches(0.6),
                 "▼  读写", font_size=18, font_color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)
    
    # Layer C: 数据层
    add_rect(slide, Inches(1.5), Inches(8.1), Inches(25), Inches(1.8),
             RGBColor(0x14, 0x1A, 0x2E), border_color=RGBColor(0x6B, 0x7B, 0xFF))
    add_text_box(slide, Inches(2), Inches(8.3), Inches(4), Inches(0.6),
                 "Layer C · 数据层", font_size=24, 
                 font_color=RGBColor(0x6B, 0x7B, 0xFF), bold=True)
    add_text_box(slide, Inches(2), Inches(8.9), Inches(23), Inches(0.8),
                 "SQLite（WAL 模式）→ PostgreSQL + TimescaleDB + Redis  |  70 条知识库规则  |  行业价格数据",
                 font_size=18, font_color=SEMI_WHITE)


def create_slide5_tech_highlights(prs):
    """Slide 5: 技术亮点。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    
    add_text_box(slide, Inches(1.5), Inches(0.8), Inches(25), Inches(1.2),
                 "04  技术亮点：为什么是最强 AI 合伙人", 
                 font_size=40, font_color=TEAL, bold=True)
    
    highlights = [
        ("🧠", "三层智能路由", "CSI≤20 规则引擎（0 成本）\nCSI 21-60 Haiku（快速）\nCSI>60 Opus（深度推理）\n自动降级 fallback"),
        ("🔌", "标准 MCP 协议", "12 个工具即插即用\nFastMCP stdio/HTTP 双模式\nclawhub install 一键安装\n开源可扩展"),
        ("📡", "真实传感器接入", "涂鸦 IoT 智能水质仪\nMock/Tuya/DIMOS 三种适配器\nENV 一键切换模式\n5 分钟采集周期"),
        ("🤖", "具身智能预留", "DIMOS × 宇树机器狗\n自动巡塘路径规划\n无人机航拍水质\n投喂执行器控制"),
        ("🛡️", "工程级安全", "传感器物理校验\n危险操作人工确认\n飞书推送 60min 去重\nLLM 10s 强制超时"),
        ("📊", "知识驱动决策", "70 条养殖规则\n54 篇学术参考\nTF-IDF 检索引擎\n规则+AI 混合推理"),
    ]
    
    for i, (emoji, title, desc) in enumerate(highlights):
        row = i // 3
        col = i % 3
        x = Inches(1.5 + col * 8.5)
        y = Inches(2.5 + row * 4)
        
        add_rect(slide, x, y, Inches(7.5), Inches(3.5),
                 RGBColor(0x14, 0x1A, 0x2E), border_color=TEAL)
        
        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(6.5), Inches(0.8),
                     f"{emoji}  {title}", font_size=26, font_color=WHITE, bold=True)
        
        add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(6.5), Inches(2),
                     desc, font_size=18, font_color=SEMI_WHITE)


def create_slide6_demo(prs):
    """Slide 6: 效果展示。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    
    add_text_box(slide, Inches(1.5), Inches(0.8), Inches(25), Inches(1.2),
                 "05  效果展示：从传感器到飞书，全链路自动化", 
                 font_size=40, font_color=TEAL, bold=True)
    
    # 流程步骤
    steps = [
        ("1️⃣", "传感器采集", "水温 26.5°C\nDO 3.2 mg/L ⚠️\npH 7.8\n氨氮 0.15 mg/L"),
        ("2️⃣", "MCP 工具计算", "水质综合评分: 42\n风险等级: 中等\n关键词触发: DO<3.0"),
        ("3️⃣", "Agent 决策", "模型: Haiku\nrisk_level: 3\n建议: 立即开增氧机\n减少投喂 50%"),
        ("4️⃣", "飞书推送", "🚨 虾塘 A3 告警\n溶解氧低于安全阈值\n已自动触发增氧\n技术员已收到通知"),
    ]
    
    for i, (num, title, content) in enumerate(steps):
        x = Inches(1.2 + i * 6.5)
        add_rect(slide, x, Inches(2.8), Inches(5.8), Inches(6.5),
                 RGBColor(0x14, 0x1A, 0x2E), border_color=TEAL)
        
        add_text_box(slide, x + Inches(0.3), Inches(3.0), Inches(5), Inches(0.6),
                     f"{num} {title}", font_size=28, font_color=TEAL, bold=True)
        
        add_text_box(slide, x + Inches(0.3), Inches(3.8), Inches(5), Inches(4.5),
                     content, font_size=22, font_color=SEMI_WHITE)
        
        # 箭头（除最后一个）
        if i < 3:
            add_text_box(slide, x + Inches(5.8), Inches(5.5), Inches(0.7), Inches(0.6),
                         "→", font_size=36, font_color=TEAL, bold=True,
                         alignment=PP_ALIGN.CENTER)


def create_slide7_business(prs):
    """Slide 7: 商业价值。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    
    add_text_box(slide, Inches(1.5), Inches(0.8), Inches(25), Inches(1.2),
                 "06  商业价值：替代 3 个岗位，效率提升 10 倍", 
                 font_size=40, font_color=TEAL, bold=True)
    
    # 左侧：大数字
    metrics = [
        ("3,050 万亩", "中国虾塘总面积", TEAL),
        ("95%", "无专业 AI 覆盖率", ACCENT_ORANGE),
        ("¥1,800/月", "SaaS 单塘月费", WHITE),
        ("1.5%", "撮合交易佣金", WHITE),
    ]
    
    for i, (number, label, color) in enumerate(metrics):
        y = Inches(2.5 + i * 1.8)
        add_text_box(slide, Inches(1.5), y, Inches(6), Inches(0.8),
                     number, font_size=48, font_color=color, bold=True)
        add_text_box(slide, Inches(1.5), y + Inches(0.8), Inches(6), Inches(0.6),
                     label, font_size=20, font_color=LIGHT_GRAY)
    
    # 右侧：替代岗位
    add_rect(slide, Inches(9), Inches(2.5), Inches(16.5), Inches(7),
             RGBColor(0x14, 0x1A, 0x2E), border_color=TEAL)
    
    add_text_box(slide, Inches(9.5), Inches(2.8), Inches(15), Inches(0.8),
                 "AI 替代的岗位", font_size=28, font_color=TEAL, bold=True)
    
    roles = [
        ("👤 24h 值班技术员", "→ 哨兵 Agent 全天候自动监控", "年省 ¥8-12 万/塘"),
        ("👤 养殖顾问", "→ 策略 Agent 每日精准决策", "年省 ¥5-8 万/塘"),
        ("👤 销售经理 + 市场分析师", "→ 增长 Agent 自动获客+市场分析", "年省 ¥10-15 万/塘"),
    ]
    
    for i, (role, replacement, saving) in enumerate(roles):
        y = Inches(3.8 + i * 1.8)
        add_text_box(slide, Inches(9.5), y, Inches(7), Inches(0.6),
                     role, font_size=22, font_color=WHITE, bold=True)
        add_text_box(slide, Inches(9.5), y + Inches(0.6), Inches(7), Inches(0.5),
                     replacement, font_size=18, font_color=SEMI_WHITE)
        add_text_box(slide, Inches(17), y + Inches(0.2), Inches(7), Inches(0.6),
                     saving, font_size=22, font_color=ACCENT_ORANGE, bold=True)


def create_slide8_closing(prs):
    """Slide 8: 结尾。"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, DARK_BG)
    
    add_text_box(slide, Inches(2), Inches(2), Inches(24), Inches(2),
                 "🦞 虾塘大亨", font_size=72, font_color=TEAL, bold=True,
                 alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2), Inches(4.2), Inches(24), Inches(1),
                 "让每一口虾塘，都有一个 AI 合伙人。", font_size=36, 
                 font_color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2), Inches(6), Inches(24), Inches(1),
                 "开源  ·  即插即用  ·  标准 MCP 协议", font_size=24,
                 font_color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2), Inches(7.5), Inches(24), Inches(0.8),
                 "GitHub: github.com/Alex647648/shrimp-tycoon", font_size=20,
                 font_color=SEMI_WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2), Inches(8.5), Inches(24), Inches(0.8),
                 "谢谢！", font_size=48, font_color=TEAL, bold=True,
                 alignment=PP_ALIGN.CENTER)


def main():
    # 创建新 PPT，使用模板尺寸（28" x 10.5"）
    prs = Presentation()
    prs.slide_width = Inches(28)
    prs.slide_height = Inches(10.5)
    
    # 创建 8 页路演
    create_slide1_cover(prs)      # 封面
    create_slide2_problem(prs)    # 行业痛点
    create_slide3_solution(prs)   # 解决方案
    create_slide4_architecture(prs)  # 技术架构
    create_slide5_tech_highlights(prs)  # 技术亮点
    create_slide6_demo(prs)       # 效果展示
    create_slide7_business(prs)   # 商业价值
    create_slide8_closing(prs)    # 结尾
    
    prs.save(OUTPUT)
    print(f"✅ PPT 已生成：{OUTPUT}")
    print(f"   总页数：{len(prs.slides)}")


if __name__ == "__main__":
    main()
